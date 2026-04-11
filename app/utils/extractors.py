from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass

from werkzeug.datastructures import FileStorage


@dataclass(slots=True)
class ExtractResult:
    filename: str
    content_type: str | None
    ok: bool
    body: list[str]
    error: str | None = None


def _ext(filename: str) -> str:
    return os.path.splitext(filename)[1].lower().lstrip(".")


def _extract_pdf_pages(data: bytes) -> list[str]:
    try:
        import pdfplumber  # type: ignore[import-not-found]
    except Exception as e:
        raise RuntimeError("dependency_missing: pdfplumber") from e

    out: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            txt = page.extract_text() or ""
            out.append(txt.strip())
    return out


def _extract_pptx_slides(data: bytes) -> list[str]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except Exception as e:
        raise RuntimeError("dependency_missing: python-pptx") from e

    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for slide in prs.slides:
        slide_lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                slide_lines.append(text.strip())
        out.append("\n".join(slide_lines).strip())
    return out


def _convert_office_to_pdf(data: bytes, filename: str) -> bytes:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("dependency_missing: libreoffice")

    ext = _ext(filename) or "ppt"

    with tempfile.TemporaryDirectory(prefix="vibe-office-") as tmpdir:
        in_path = os.path.join(tmpdir, f"input.{ext}")
        out_dir = os.path.join(tmpdir, "out")
        os.makedirs(out_dir, exist_ok=True)

        with open(in_path, "wb") as f:
            f.write(data)

        # libreoffice writes output to out_dir with same basename and .pdf extension
        cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--nodefault",
            "--norestore",
            "--convert-to",
            "pdf",
            "--outdir",
            out_dir,
            in_path,
        ]
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError as e:
            stderr = (e.stderr or b"").decode("utf-8", errors="replace").strip()
            raise RuntimeError(f"office_to_pdf_failed: {stderr or 'unknown_error'}") from e

        base = os.path.splitext(os.path.basename(in_path))[0]
        pdf_path = os.path.join(out_dir, f"{base}.pdf")
        if not os.path.exists(pdf_path):
            candidates = [p for p in os.listdir(out_dir) if p.lower().endswith(".pdf")]
            raise RuntimeError(f"office_to_pdf_failed: pdf_not_found ({', '.join(candidates) or 'no_outputs'})")

        with open(pdf_path, "rb") as f:
            return f.read()


def extract_text_from_upload(upload: FileStorage) -> ExtractResult:
    filename = upload.filename or "upload"
    content_type = upload.mimetype

    try:
        data = upload.read()
        if not data:
            return ExtractResult(
                filename=filename,
                content_type=content_type,
                ok=False,
                body=[],
                error="empty_file",
            )

        ext = _ext(filename)
        if ext == "pdf":
            body = _extract_pdf_pages(data)
        elif ext in {"pptx", "ppt"}:
            if ext == "ppt":
                pdf_bytes = _convert_office_to_pdf(data, filename)
                body = _extract_pdf_pages(pdf_bytes)
            else:
                body = _extract_pptx_slides(data)
        else:
            return ExtractResult(
                filename=filename,
                content_type=content_type,
                ok=False,
                body=[],
                error="unsupported_extension",
            )

        return ExtractResult(
            filename=filename,
            content_type=content_type,
            ok=True,
            body=body,
            error=None,
        )
    except Exception as e:  # noqa: BLE001 - boundary catches for API response
        return ExtractResult(
            filename=filename,
            content_type=content_type,
            ok=False,
            body=[],
            error=f"{type(e).__name__}: {e}",
        )

